"""Restructure immunoverse_a.pptx per mentor feedback.

New slide order:
  0  Title — updated tagline
  1  The Efficiency Problem (brief)
  2  The Fragmentation Reality ("Great tools don't talk")
  3  NEW — Three Researchers composability hook
  4  The ImmunoVerse Approach (updated: 7 agents)
  5  The 4-Stage Action Gauntlet
  6  Why "Composable" is Better
  7  The Economic Shift
  8  Handling Data Realities (Bias)
  9  The Composable Primitive (updated close)

Dropped: Mandatory Engine Test, Project Readiness, Image Sources
"""

import copy
import sys
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor

# Brand colours (sampled from existing slides)
DARK_NAVY   = RGBColor(30, 41, 59)
MED_SLATE   = RGBColor(71, 85, 105)
LIGHT_SLATE = RGBColor(100, 116, 139)

NS_A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P   = "http://schemas.openxmlformats.org/presentationml/2006/main"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_run(run, text, size_pt, bold=False, color=DARK_NAVY):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_tb_single(tb, text, size_pt, bold=False, color=DARK_NAVY):
    """Overwrite the first paragraph of a text frame with a single run."""
    tf = tb.text_frame
    tf.word_wrap = True
    # Clear all paragraphs
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    # Remove all runs
    for r in p.runs:
        r._r.getparent().remove(r._r)
    run = p.add_run()
    _set_run(run, text, size_pt, bold, color)


def _add_paragraph(tf, text, size_pt=10, bold=False, color=MED_SLATE, bullet=False):
    """Append a paragraph to an existing text frame."""
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    run = p.add_run()
    _set_run(run, text, size_pt, bold, color)
    return p


def move_slide(prs, old_idx, new_idx):
    lst = prs.slides._sldIdLst
    items = list(lst)
    el = items[old_idx]
    lst.remove(el)
    lst.insert(new_idx, el)


def remove_slide(prs, idx):
    lst = prs.slides._sldIdLst
    items = list(lst)
    lst.remove(items[idx])


def duplicate_slide(prs, template_idx):
    """Duplicate a slide (including image relationships) and append at end."""
    template = prs.slides[template_idx]
    layout   = template.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Build rId remapping for image relationships
    rId_map = {}
    for rId, rel in template.part.rels.items():
        if "image" in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    # Replace spTree content
    new_tree  = new_slide.shapes._spTree
    tmpl_tree = template.shapes._spTree

    for child in list(new_tree):
        new_tree.remove(child)

    for child in tmpl_tree:
        node = copy.deepcopy(child)
        # Patch r:embed in <a:blip> for images
        for blip in node.iter(f"{{{NS_A}}}blip"):
            old = blip.get(f"{{{NS_R}}}embed")
            if old and old in rId_map:
                blip.set(f"{{{NS_R}}}embed", rId_map[old])
        new_tree.append(node)

    return len(prs.slides) - 1


def find_text_shapes(slide):
    """Return list of (shape, first_para_text) for text-bearing shapes."""
    result = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                result.append((shape, t))
    return result


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    src  = "slides/immunoverse_a.pptx"
    dest = "slides/immunoverse_pitch.pptx"

    prs = Presentation(src)

    # -----------------------------------------------------------------------
    # Step 1 — Update title slide tagline (slide index 0)
    # -----------------------------------------------------------------------
    title_slide = prs.slides[0]
    for shape in title_slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "decentralized agent network" in t:
                _set_tb_single(
                    shape,
                    "Gene therapy safety screening as a composable network service.",
                    size_pt=24, bold=False, color=LIGHT_SLATE,
                )
                break

    # -----------------------------------------------------------------------
    # Step 2 — Create "Three Researchers" slide by duplicating slide 2
    #          (Fragmentation Reality) and rewriting all text content
    # -----------------------------------------------------------------------
    new_idx = duplicate_slide(prs, 2)   # appended at end
    three_slide = prs.slides[new_idx]

    # Identify and rewrite each text shape by its current content
    for shape in three_slide.shapes:
        if not shape.has_text_frame:
            continue
        original = shape.text_frame.text.strip()

        # Page header (large title at very top)
        if "Fragmentation Reality" in original:
            _set_tb_single(shape, "The Composability Moment",
                           size_pt=29, bold=True, color=DARK_NAVY)

        # Sub-heading on the left: "Great tools, but they don't talk."
        elif "Great tools" in original:
            _set_tb_single(shape, "Three Researchers. One Message.",
                           size_pt=20, bold=True, color=DARK_NAVY)

        # Left body — "The industry already has powerful…"
        elif "industry already has" in original:
            tf = shape.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in p0.runs:
                r._r.getparent().remove(r._r)
            r0 = p0.add_run()
            _set_run(r0, "Before ImmunoVerse:", size_pt=10, bold=False, color=MED_SLATE)

        # Bottom left callout: "When data is siloed…"
        elif "When data is siloed" in original:
            tf = shape.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in p0.runs:
                r._r.getparent().remove(r._r)
            r0 = p0.add_run()
            _set_run(r0,
                     "Dr. Chen's 3 years of HLA expertise — now a callable network service. "
                     "No integration. No file formats. Just an address.",
                     size_pt=10, bold=False, color=MED_SLATE)

        # Bullet "Idiosyncratic file formats"
        elif "Idiosyncratic file formats" in original:
            _set_tb_single(shape,
                           "Dr. Chen: IEDB HLA binding — manual CSV exports, 2 days per candidate",
                           size_pt=10, bold=False, color=MED_SLATE)

        # Bullet "Scattered, unlinked databases"
        elif "Scattered, unlinked databases" in original:
            _set_tb_single(shape,
                           "Dr. Kim: NetTCR T-cell models — Jupyter notebook, no public API",
                           size_pt=10, bold=False, color=MED_SLATE)

        # Bullet "One-off retrieval scripts"
        elif "One-off retrieval scripts" in original:
            _set_tb_single(shape,
                           "Dr. Patel: needs both, right now — emails, waits days, reformats data",
                           size_pt=10, bold=False, color=MED_SLATE)

        # Right panel heading: "The 'Old City' Problem"
        elif "Old City" in original:
            _set_tb_single(shape, "With ImmunoVerse on Agentverse:",
                           size_pt=20, bold=True, color=DARK_NAVY)

        # Right panel body
        elif "Navigating bio-data today" in original:
            tf = shape.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in p0.runs:
                r._r.getparent().remove(r._r)
            r0 = p0.add_run()
            _set_run(r0,
                     "Dr. Patel sends ONE message to one Agentverse address. "
                     "All 4 stages run in parallel. Results in seconds.\n\n"
                     "→  Dr. Chen's HLA work: a callable service\n"
                     "→  Dr. Kim's TCR work: a callable service\n"
                     "→  Stage 2 is live on Agentverse right now",
                     size_pt=10, bold=False, color=MED_SLATE)

    # -----------------------------------------------------------------------
    # Step 3 — Move the new slide to position 3
    #          (after Title, Efficiency Problem, Fragmentation Reality)
    # -----------------------------------------------------------------------
    #  Before move: slides 0-12 (original 0-11 + new at 12)
    move_slide(prs, new_idx, 3)

    # After this move the layout is:
    #  0  Title
    #  1  Efficiency Problem
    #  2  Fragmentation Reality
    #  3  Three Researchers  ← inserted
    #  4  ImmunoVerse Approach   (was 3)
    #  5  4-Stage Gauntlet       (was 4)
    #  6  Why Composable         (was 5)
    #  7  Economic Shift         (was 6)
    #  8  Mandatory Engine Test  (was 7)
    #  9  Bias                   (was 8)
    #  10 Project Readiness      (was 9)
    #  11 Composable Primitive   (was 10)
    #  12 Image Sources          (was 11)

    # -----------------------------------------------------------------------
    # Step 4 — Update "The ImmunoVerse Approach" slide (now index 4)
    # -----------------------------------------------------------------------
    approach_slide = prs.slides[4]
    for shape in approach_slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if "The ImmunoVerse Approach" in t:
            _set_tb_single(shape, "7 Agents. Any One Callable.",
                           size_pt=29, bold=True, color=DARK_NAVY)

        elif "A Team of Autonomous Agents" in t:
            _set_tb_single(shape, "A Composable Agent Network",
                           size_pt=20, bold=True, color=DARK_NAVY)

        elif "Built on the" in t or "Fetch.ai Agentverse" in t:
            tf = shape.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in p0.runs:
                r._r.getparent().remove(r._r)
            r0 = p0.add_run()
            _set_run(r0,
                     "7 specialist agents deployed on Fetch.ai Agentverse. "
                     "Each is independently addressable — any researcher or system can call "
                     "a single stage without running the full pipeline.",
                     size_pt=10, bold=False, color=MED_SLATE)

        elif "Each agent acts as a specialist" in t:
            tf = shape.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in p0.runs:
                r._r.getparent().remove(r._r)
            r0 = p0.add_run()
            _set_run(r0,
                     "Orchestrator → Stage 1 (ESMFold) → Stage 2 (HLA/IEDB) → "
                     "Stage 3a (NetTCR) + Stage 3b (BepiPred) → Stage 4 (GenBio AIDO) → Report\n\n"
                     "Or skip the orchestrator. Call Stage 2 directly. Your choice.",
                     size_pt=10, bold=False, color=MED_SLATE)

    # -----------------------------------------------------------------------
    # Step 5 — Update "The Composable Primitive" closing slide (now index 11)
    # -----------------------------------------------------------------------
    close_slide = prs.slides[11]
    for shape in close_slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if "Moving safety from" in t:
            _set_tb_single(shape,
                           "Stage 2 is a standalone HLA binding service on Agentverse right now.\n"
                           "Send it a sequence. Get a binding prediction back in seconds.",
                           size_pt=19, bold=False, color=MED_SLATE)

        elif "QUESTIONS?" in t:
            _set_tb_single(shape,
                           "QUESTIONS?  HERE'S THE AGENTVERSE ADDRESS.",
                           size_pt=16, bold=False, color=LIGHT_SLATE)

    # -----------------------------------------------------------------------
    # Step 6 — Remove unwanted slides (highest index first to avoid shifting)
    #          8: Mandatory Engine Test
    #          10: Project Readiness
    #          12: Image Sources
    # -----------------------------------------------------------------------
    remove_slide(prs, 12)   # Image Sources
    remove_slide(prs, 10)   # Project Readiness
    remove_slide(prs, 8)    # Mandatory Engine Test

    # Final order after removals:
    #  0  Title
    #  1  Efficiency Problem
    #  2  Fragmentation Reality
    #  3  Three Researchers
    #  4  ImmunoVerse Approach
    #  5  4-Stage Gauntlet
    #  6  Why Composable
    #  7  Economic Shift
    #  8  Bias
    #  9  Composable Primitive

    prs.save(dest)
    print(f"Saved {dest} ({len(prs.slides)} slides)")
    for i, slide in enumerate(prs.slides):
        titles = [s.text_frame.text.strip()
                  for s in slide.shapes
                  if s.has_text_frame and s.text_frame.text.strip()]
        headline = titles[0] if titles else "(no text)"
        print(f"  {i+1:2d}. {headline[:70]}")


if __name__ == "__main__":
    main()
